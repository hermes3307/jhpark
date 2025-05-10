import os
import time
import json
import hashlib
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any
import logging

# Document processing libraries
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import pandas as pd

# ML/Vector libraries
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma, FAISS
from langchain.document_loaders import TextLoader
import chromadb

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_processing.log'),
        logging.StreamHandler()
    ]
)

class DocumentVectorizer:
    """Main class for processing and storing office documents in vector database"""
    
    def __init__(self, batch_folder: str, vector_db_path: str, use_openai: bool = False):
        """
        Initialize the document vectorizer
        
        Args:
            batch_folder: Path to folder containing documents to process
            vector_db_path: Path to store vector database
            use_openai: Whether to use OpenAI embeddings (requires API key)
        """
        self.batch_folder = Path(batch_folder)
        self.vector_db_path = Path(vector_db_path)
        self.processed_files_log = self.vector_db_path / "processed_files.json"
        
        # Create directories if they don't exist
        self.batch_folder.mkdir(parents=True, exist_ok=True)
        self.vector_db_path.mkdir(parents=True, exist_ok=True)
        
        # Initialize embeddings
        if use_openai:
            self.embeddings = OpenAIEmbeddings()
        else:
            self.embeddings = HuggingFaceEmbeddings(
                model_name="all-MiniLM-L6-v2"
            )
        
        # Initialize vector store
        self.vector_store = self._initialize_vector_store()
        
        # Load processed files log
        self.processed_files = self._load_processed_files()
        
        # Supported file extensions by category
        self.file_categories = {
            'documents': ['.pdf', '.doc', '.docx', '.txt', '.rtf'],
            'spreadsheets': ['.xlsx', '.xls', '.csv'],
            'presentations': ['.pptx', '.ppt'],
            'text': ['.txt', '.md', '.json', '.yaml', '.yml']
        }
        
    def _initialize_vector_store(self):
        """Initialize or load existing vector store"""
        try:
            # Try to load existing FAISS index
            if (self.vector_db_path / "index.faiss").exists():
                return FAISS.load_local(
                    str(self.vector_db_path), 
                    self.embeddings
                )
            else:
                # Create new FAISS index
                return FAISS.from_texts([""], self.embeddings)
        except Exception as e:
            logging.error(f"Error initializing vector store: {e}")
            return FAISS.from_texts([""], self.embeddings)
    
    def _load_processed_files(self) -> Dict[str, Any]:
        """Load log of processed files"""
        if self.processed_files_log.exists():
            with open(self.processed_files_log, 'r') as f:
                return json.load(f)
        return {}
    
    def _save_processed_files(self):
        """Save log of processed files"""
        with open(self.processed_files_log, 'w') as f:
            json.dump(self.processed_files, f, indent=2)
    
    def _get_file_hash(self, file_path: Path) -> str:
        """Generate hash of file content"""
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            buf = f.read()
            hasher.update(buf)
        return hasher.hexdigest()
    
    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            logging.error(f"Error extracting text from PDF {file_path}: {e}")
        return text
    
    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logging.error(f"Error extracting text from DOCX {file_path}: {e}")
            return ""
    
    def _extract_text_from_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel file"""
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            text = ""
            for sheet_name, sheet_df in df.items():
                text += f"Sheet: {sheet_name}\n"
                text += sheet_df.to_string() + "\n\n"
            return text
        except Exception as e:
            logging.error(f"Error extracting text from Excel {file_path}: {e}")
            return ""
    
    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PowerPoint {file_path}: {e}")
            return ""
    
    def _extract_text_from_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logging.error(f"Error extracting text from CSV {file_path}: {e}")
            return ""
    
    def _extract_text_from_txt(self, file_path: Path) -> str:
        """Extract text from text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logging.error(f"Error extracting text from TXT {file_path}: {e}")
            return ""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from any supported file type"""
        ext = file_path.suffix.lower()
        
        extraction_methods = {
            '.pdf': self._extract_text_from_pdf,
            '.docx': self._extract_text_from_docx,
            '.doc': self._extract_text_from_docx,
            '.xlsx': self._extract_text_from_xlsx,
            '.xls': self._extract_text_from_xlsx,
            '.pptx': self._extract_text_from_pptx,
            '.ppt': self._extract_text_from_pptx,
            '.csv': self._extract_text_from_csv,
            '.txt': self._extract_text_from_txt,
            '.md': self._extract_text_from_txt,
            '.json': self._extract_text_from_txt,
            '.yaml': self._extract_text_from_txt,
            '.yml': self._extract_text_from_txt,
        }
        
        if ext in extraction_methods:
            return extraction_methods[ext](file_path)
        else:
            logging.warning(f"Unsupported file type: {ext}")
            return ""
    
    def get_file_category(self, file_path: Path) -> str:
        """Determine the category of a file based on its extension"""
        ext = file_path.suffix.lower()
        for category, extensions in self.file_categories.items():
            if ext in extensions:
                return category
        return "other"
    
    def process_document(self, file_path: Path) -> bool:
        """Process a single document and add to vector database"""
        try:
            # Check if file has already been processed
            file_hash = self._get_file_hash(file_path)
            if str(file_path) in self.processed_files:
                if self.processed_files[str(file_path)]['hash'] == file_hash:
                    logging.info(f"File already processed: {file_path}")
                    return False
            
            # Extract text from document
            text = self.extract_text(file_path)
            if not text.strip():
                logging.warning(f"No text extracted from {file_path}")
                return False
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
            )
            chunks = text_splitter.split_text(text)
            
            # Create metadata for each chunk
            category = self.get_file_category(file_path)
            metadata = {
                'source': str(file_path),
                'filename': file_path.name,
                'category': category,
                'processed_date': datetime.now().isoformat(),
                'file_hash': file_hash
            }
            
            # Add chunks to vector store
            texts_with_metadata = [(chunk, metadata) for chunk in chunks]
            self.vector_store.add_texts(
                texts=[t[0] for t in texts_with_metadata],
                metadatas=[t[1] for t in texts_with_metadata]
            )
            
            # Update processed files log
            self.processed_files[str(file_path)] = {
                'hash': file_hash,
                'processed_date': datetime.now().isoformat(),
                'category': category,
                'chunks': len(chunks)
            }
            self._save_processed_files()
            
            logging.info(f"Successfully processed {file_path} ({len(chunks)} chunks)")
            return True
            
        except Exception as e:
            logging.error(f"Error processing document {file_path}: {e}")
            return False
    
    def batch_process_folder(self) -> Dict[str, int]:
        """Process all new files in the batch folder"""
        stats = {
            'total_files': 0,
            'processed': 0,
            'skipped': 0,
            'errors': 0
        }
        
        # Get all supported files in batch folder
        supported_extensions = set()
        for extensions in self.file_categories.values():
            supported_extensions.update(extensions)
        
        for file_path in self.batch_folder.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                stats['total_files'] += 1
                
                if self.process_document(file_path):
                    stats['processed'] += 1
                else:
                    stats['skipped'] += 1
        
        # Save vector store
        self.vector_store.save_local(str(self.vector_db_path))
        
        logging.info(f"Batch processing complete: {stats}")
        return stats
    
    def search_documents(self, query: str, k: int = 5) -> List[Dict[str, Any]]:
        """Search documents in the vector database"""
        try:
            docs = self.vector_store.similarity_search_with_score(query, k=k)
            results = []
            for doc, score in docs:
                results.append({
                    'content': doc.page_content,
                    'metadata': doc.metadata,
                    'score': score
                })
            return results
        except Exception as e:
            logging.error(f"Error searching documents: {e}")
            return []
    
    def monitor_folder(self, check_interval: int = 60):
        """Monitor batch folder for new files at regular intervals"""
        logging.info(f"Starting folder monitoring. Checking every {�_interval} seconds...")
        
        while True:
            try:
                logging.info("Checking for new files...")
                stats = self.batch_process_folder()
                
                if stats['processed'] > 0:
                    logging.info(f"Processed {stats['processed']} new files")
                else:
                    logging.info("No new files to process")
                
                time.sleep(check_interval)
                
            except KeyboardInterrupt:
                logging.info("Monitoring stopped by user")
                break
            except Exception as e:
                logging.error(f"Error during monitoring: {e}")
                time.sleep(check_interval)


# Example usage and batch script
if __name__ == "__main__":
    # Configuration
    BATCH_FOLDER = "./batch_documents"
    VECTOR_DB_PATH = "./vector_database"
    USE_OPENAI = False  # Set to True if you want to use OpenAI embeddings
    
    # Initialize the vectorizer
    vectorizer = DocumentVectorizer(
        batch_folder=BATCH_FOLDER,
        vector_db_path=VECTOR_DB_PATH,
        use_openai=USE_OPENAI
    )
    
    # Process initial batch
    print("Processing initial batch of documents...")
    stats = vectorizer.batch_process_folder()
    print(f"Initial processing complete: {stats}")
    
    # Example search
    print("\nExample search:")
    results = vectorizer.search_documents("quarterly sales report", k=3)
    for i, result in enumerate(results, 1):
        print(f"\nResult {i}:")
        print(f"File: {result['metadata']['filename']}")
        print(f"Category: {result['metadata']['category']}")
        print(f"Score: {result['score']}")
        print(f"Content: {result['content'][:200]}...")
    
    # Optional: Start monitoring for new files
    # Uncomment the following line to enable continuous monitoring
    # vectorizer.monitor_folder(check_interval=60)